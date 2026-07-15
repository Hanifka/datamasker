#!/usr/bin/env python3
"""
app.py — Streamlit UI for datamask (reversible sensitive-data masking).

Run:  streamlit run app.py   (needs datamask.py in the same directory)

Two tabs:
  🎭 Encode — scan → review/select values → tokenize → download masked + mapping.json
  🔓 Decode — restore originals in AI output (pptx/docx/xlsx/csv/text) via mapping.json

Results persist across reruns (session_state), so download buttons don't vanish
after the first click.
"""

import io
import json
import os
import re
import tempfile
import zipfile

import streamlit as st

# ==== inlined from datamask.py ====
import csv as csv_mod
import io
import json
import os
import re
import sys
from collections import OrderedDict
from datetime import datetime, timezone

# ---------------------------------------------------------------- detectors

PRIVATE_NETS = re.compile(
    r"^(10\.|192\.168\.|172\.(1[6-9]|2[0-9]|3[01])\.|127\.|169\.254\.)"
)

def _valid_ipv4(s: str) -> bool:
    parts = s.split(".")
    return len(parts) == 4 and all(p.isdigit() and 0 <= int(p) <= 255 for p in parts)

# Order matters: more specific first (email before domain, FQDN before hostname).
DETECTORS = [
    # (category, token prefix, regex, post-filter)
    ("email",    "MAILX", re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"), None),
    ("ipv4",     "IPX",   re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b"), _valid_ipv4),
    ("ipv6",     "IP6X",  re.compile(r"\b(?:[0-9A-Fa-f]{1,4}:){2,7}[0-9A-Fa-f]{1,4}\b"),
        lambda s: (re.search(r"[A-Fa-f]", s) is not None                      # needs hex letter
                   and not re.fullmatch(r"(?:[0-9A-Fa-f]{2}:){5}[0-9A-Fa-f]{2}", s)  # not a MAC
                   and not re.fullmatch(r"\d{1,2}:\d{2}(:\d{2})?", s))),      # not a timestamp
    ("mac",      "MACX",  re.compile(r"\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b"), None),
    # Windows-style hostnames: DESKTOP-XXXX, SRV-*, WIN-*, DC01, LAPTOP-...
    ("hostname", "HOSTX", re.compile(
        r"\b(?:DESKTOP|LAPTOP|WIN|SRV|SVR|DC|PC|WS|HOST|VM|APP|DB|WEB|MAIL|FW|SW|RTR|PRD|DEV|UAT)"
        r"[-_][A-Za-z0-9][A-Za-z0-9-_]{1,30}\b", re.IGNORECASE), None),
    # FQDN (internal domains etc.) — post-filter drops common file extensions
    ("fqdn",     "FQDNX", re.compile(
        r"\b(?:[A-Za-z0-9](?:[A-Za-z0-9-]{0,61}[A-Za-z0-9])?\.){2,}[A-Za-z]{2,}\b"),
        lambda s: not re.search(r"\.(?:exe|dll|csv|xlsx|docx|pptx|txt|log|json|xml|py|sh|js|zip|png|jpg)$", s, re.I)),
    # log key=value hostnames: host=WEB01, hostname: jakarta-app02, computer="PC-BUDI"
    ("hostname", "HOSTX", re.compile(
        r"(?i)\b(?:src[._-]?host|dst[._-]?host|host(?:name)?|computer(?:name)?|"
        r"device(?:name)?|machine|agent[._-]?name|asset|node)\s*[:=]\s*\"?"
        r"([A-Za-z][A-Za-z0-9._-]{2,60})\"?"),
        lambda s: s.lower() not in {"true","false","null","none","unknown","localhost",
                                    "name","value","string","header","info"}),
    # DOMAIN\username — filter out Windows paths / registry keys / filenames
    ("username", "USERX", re.compile(r"\b[A-Za-z][A-Za-z0-9_-]{1,20}\\[A-Za-z][A-Za-z0-9._-]{2,30}\b"),
        lambda s: (lambda dom, usr: (
            dom.lower() not in {
                "downloads","documents","desktop","users","user","windows","system32",
                "program","programs","programdata","appdata","temp","tmp","local","roaming",
                "currentversion","software","microsoft","policies","services","system",
                "wow6432node","classes","run","runonce","uninstall","explorer","shell",
                "drivers","etc","config","control","setup","installer","packages","start",
                "menu","common","public","default","all"}
            and usr.lower() not in {
                "uninstall","currentversion","system32","run","runonce","explorer",
                "shell","policies","parameters","services","software","windows"}
            and not re.search(r"\.(exe|dll|sys|msi|bat|cmd|ps1|vbs|js|zip|rar|7z|iso|"
                              r"pdf|docx?|xlsx?|pptx?|txt|log|csv|json|xml|ini|dat|tmp|"
                              r"lnk|url|jpg|jpeg|png|gif|bmp)\.?$", usr, re.I)
        ))(*s.split("\\", 1))),
    # key=value contextual usernames: user=jdoe, username: budi.s, account="ops_admin"
    ("username", "USERX", re.compile(
        r"(?i)\b(?:user(?:name)?|account|login|logon|acct)\s*[:=]\s*\"?([A-Za-z][A-Za-z0-9._-]{2,30})\"?"), None),
]

MASK_RECOMMENDATION = {
    "ipv4":     "Tokenize (reversible). If sharing externally without codebook: mask last octet (10.1.2.x).",
    "ipv6":     "Tokenize (reversible). External: truncate to /48 prefix.",
    "hostname": "Tokenize (reversible). External: role-based alias (WEB-SRV-A).",
    "fqdn":     "Tokenize (reversible). External: replace internal domain with example.internal.",
    "username": "Tokenize (reversible). Never share real usernames outside the org — high phishing value.",
    "email":    "Tokenize (reversible). External: hash local-part, keep domain if public.",
    "mac":      "Tokenize (reversible). External: keep OUI (vendor) bytes, mask last 3 octets.",
}

# ---------------------------------------------------------------- codebook

class Codebook:
    def __init__(self, path=None):
        self.map = OrderedDict()      # original -> token
        self.rev = OrderedDict()      # token -> original
        self.meta = {"created": datetime.now(timezone.utc).isoformat(), "tool": "datamask v1.0"}
        self.counters = {}
        if path and os.path.exists(path):
            with open(path) as f:
                data = json.load(f)
            self.meta = data.get("meta", self.meta)
            for orig, entry in data.get("entries", {}).items():
                self.map[orig] = entry["token"]
                self.rev[entry["token"]] = orig
                self._bump_counter(entry["token"])
            self.categories = {o: e["category"] for o, e in data.get("entries", {}).items()}
        else:
            self.categories = {}

    def _bump_counter(self, token):
        m = re.match(r"([A-Z0-9]+?X)(\d+)$", token)
        if m:
            pfx, n = m.group(1), int(m.group(2))
            self.counters[pfx] = max(self.counters.get(pfx, 0), n)

    def token_for(self, original, category, prefix):
        if original in self.map:
            return self.map[original]
        self.counters[prefix] = self.counters.get(prefix, 0) + 1
        token = f"{prefix}{self.counters[prefix]:03d}"
        self.map[original] = token
        self.rev[token] = original
        self.categories[original] = category
        return token

    def save(self, path):
        entries = {o: {"token": t, "category": self.categories.get(o, "?"),
                       "recommendation": MASK_RECOMMENDATION.get(self.categories.get(o, ""), "Tokenize (reversible).")}
                   for o, t in self.map.items()}
        with open(path, "w") as f:
            json.dump({"meta": self.meta, "entries": entries}, f, indent=2)
        try:
            os.chmod(path, 0o600)
        except OSError:
            pass

# ---------------------------------------------------------------- core replace

def find_all(text):
    """Yield (category, prefix, match_text) for every detection in text."""
    out = []
    for cat, pfx, rx, flt in DETECTORS:
        for m in rx.finditer(text):
            val = m.group(1) if m.groups() else m.group(0)
            if flt and not flt(val):
                continue
            out.append((cat, pfx, val))
    return out

def mask_text(text, cb: Codebook):
    """Replace all detections in text with tokens. Longest-first to avoid partial overlap."""
    findings = find_all(text)
    if not findings:
        return text, 0
    # Assign tokens
    for cat, pfx, val in findings:
        cb.token_for(val, cat, pfx)
    # Replace longest originals first
    n = 0
    for orig in sorted({v for _, _, v in findings}, key=len, reverse=True):
        token = cb.map[orig]
        new = text.replace(orig, token)
        if new != text:
            n += text.count(orig)
            text = new
    return text, n

def unmask_text(text, cb: Codebook):
    n = 0
    for token in sorted(cb.rev, key=len, reverse=True):
        if token in text:
            n += text.count(token)
            text = text.replace(token, cb.rev[token])
    return text, n

# ---------------------------------------------------------------- file handlers

def _iter_docx_paragraphs(doc):
    """Yield every paragraph in body, tables (nested), headers, footers."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    def walk(parent):
        for p in parent.paragraphs:
            yield p
        for t in parent.tables:
            for row in t.rows:
                for cell in row.cells:
                    yield from walk(cell)
    yield from walk(doc)
    for sec in doc.sections:
        for part in (sec.header, sec.footer):
            yield from walk(part)

def _replace_in_paragraph(par, fn):
    """Apply fn(text)->(text,count) to full paragraph text, preserving run 0 formatting
    when a match spans runs. Only rewrites the paragraph if something changed."""
    full = par.text
    new, n = fn(full)
    if n == 0:
        return 0
    # Try run-by-run first (keeps formatting perfectly when matches are within a run)
    total_in_runs = 0
    for run in par.runs:
        rnew, rn = fn(run.text)
        if rn:
            run.text = rnew
            total_in_runs += rn
    if par.text == new:
        return n
    # Cross-run match remains: collapse into first run
    if par.runs:
        par.runs[0].text = new
        for run in par.runs[1:]:
            run.text = ""
    return n

def process_docx(path, out_path, fn):
    import docx
    doc = docx.Document(path)
    count = 0
    for par in _iter_docx_paragraphs(doc):
        count += _replace_in_paragraph(par, fn)
    doc.save(out_path)
    return count

def process_pptx(path, out_path, fn):
    from pptx import Presentation
    prs = Presentation(path)
    count = 0
    def handle_tf(tf):
        nonlocal count
        for par in tf.paragraphs:
            full = "".join(r.text for r in par.runs)
            new, n = fn(full)
            if n == 0:
                continue
            done_in_runs = 0
            for r in par.runs:
                rnew, rn = fn(r.text)
                if rn:
                    r.text = rnew
                    done_in_runs += rn
            cur = "".join(r.text for r in par.runs)
            if cur != new and par.runs:
                par.runs[0].text = new
                for r in par.runs[1:]:
                    r.text = ""
            count += n
    def handle_shape(shape):
        if shape.has_text_frame:
            handle_tf(shape.text_frame)
        if shape.shape_type == 6:  # group
            for s in shape.shapes:
                handle_shape(s)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    handle_tf(cell.text_frame)
    for slide in prs.slides:
        for shape in slide.shapes:
            handle_shape(shape)
        if slide.has_notes_slide:
            handle_tf(slide.notes_slide.notes_text_frame)
    prs.save(out_path)
    return count

def process_xlsx(path, out_path, fn):
    import openpyxl
    wb = openpyxl.load_workbook(path)  # formulas preserved; we only touch strings
    count = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                if isinstance(v, str) and not v.startswith("="):
                    new, n = fn(v)
                    if n:
                        cell.value = new
                        count += n
                elif isinstance(v, (int, float)) and not isinstance(v, bool):
                    # numeric cells (e.g. broken IP 172.1 stored as float)
                    s = str(v)
                    if s.endswith(".0"):
                        s2, n = fn(s[:-2])
                        if n:
                            cell.value = s2
                            count += n
                            continue
                    new, n = fn(s)
                    if n:
                        cell.value = new
                        count += n
        # sheet name itself can contain hostnames
        new_title, n = fn(ws.title)
        if n:
            ws.title = new_title[:31]
            count += n
    wb.save(out_path)
    return count

def process_textlike(path, out_path, fn):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    new, n = fn(text)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(new)
    return n

HANDLERS = {
    ".docx": process_docx,
    ".pptx": process_pptx,
    ".xlsx": process_xlsx,
    ".xlsm": process_xlsx,
    ".csv":  process_textlike,
    ".tsv":  process_textlike,
    ".txt":  process_textlike,
    ".log":  process_textlike,
    ".json": process_textlike,
    ".xml":  process_textlike,
    ".md":   process_textlike,
}

def extract_text(path):
    """Read all visible text from a file for scan-only mode."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        import docx
        doc = docx.Document(path)
        return "\n".join(p.text for p in _iter_docx_paragraphs(doc))
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            chunks.append(cell.text)
        return "\n".join(chunks)
    if ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True)
        chunks = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                chunks.extend(str(v) for v in row if isinstance(v, str))
        return "\n".join(chunks)
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ------------------------------------------------- column-aware detection

VALUE_STOPLIST = {"n/a","na","-","--","none","null","unknown","true","false","yes","no","total"}

AGGREGATE_TOKENS = {"count","counts","total","totals","sum","avg","average","mean",
                    "unique","distinct","pct","percent","percentage","qty","quantity",
                    "min","max","rate","ratio","freq","frequency"}

def header_category(header):
    """Map a column header to a sensitive category, or None."""
    if header is None:
        return None
    tokens = [t for t in re.split(r"[^a-z0-9]+", str(header).lower()) if t]
    joined = "".join(tokens)
    if any(t in AGGREGATE_TOKENS for t in tokens):
        return None  # "Source IP (Unique Count)" = kolom agregat, bukan data
    def has(pred):
        return any(pred(t) for t in tokens) or pred(joined)
    if has(lambda t: t in {"host","hostname","computer","computername","device",
                           "devicename","machine","server","servername","agent",
                           "agentname","asset","node","endpoint","workstation"}):
        return "hostname"
    if has(lambda t: t == "ip" or t in {"ipaddress","ipaddr"}
                     or (t.endswith("ip") and t not in {"zip","tooltip","ownership"})):
        return "ipv4"
    if has(lambda t: t in {"user","username","account","login","logon","acct","userid",
                           "accountname","targetuser","subjectuser"}):
        return "username"
    if has(lambda t: t in {"email","mail","emailaddress","sender","recipient"}):
        return "email"
    if has(lambda t: t in {"mac","macaddress","macaddr"}):
        return "mac"
    if has(lambda t: t in {"domain","fqdn","dns","dnsname"}):
        return "fqdn"
    return None

def find_structured(path):
    """Column-header-aware findings for spreadsheets: [(category, value, column)]."""
    ext = os.path.splitext(path)[1].lower()
    out = []
    def scan_grid(rows_iter):
        headers = None
        raw_headers = None
        for row in rows_iter:
            if headers is None:
                raw_headers = [str(c) if c is not None else "" for c in row]
                headers = [header_category(c) for c in row]
                continue
            for i, v in enumerate(row):
                if i >= len(headers) or not headers[i] or v in (None, ""):
                    continue
                s = str(v).strip()
                if (len(s) < 3 or s.lower() in VALUE_STOPLIST
                        or s.startswith("=") or (s.isdigit() and len(s) < 7)):
                    continue
                cat = headers[i]
                # value must look plausible for its category
                if cat == "ipv4" and not re.fullmatch(r"[\d.,:/\s*x-]+", s):
                    continue  # "Multiple (1,021)" etc.
                if cat == "mac" and not re.fullmatch(r"[0-9A-Fa-f:.-]+", s):
                    continue
                if cat == "email" and "@" not in s:
                    continue
                out.append((cat, s, raw_headers[i]))
    if ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True)
        for ws in wb.worksheets:
            scan_grid(ws.iter_rows(values_only=True))
    elif ext in (".csv", ".tsv"):
        delim = "\t" if ext == ".tsv" else ","
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            scan_grid(csv_mod.reader(f, delimiter=delim))
    return out


# ==== end inline ====


st.set_page_config(page_title="DataMask", page_icon="🛡️", layout="wide")
st.title("🛡️ DataMask")
st.caption("Encode data sensitif → kasih ke AI → Decode hasilnya. "
           "mapping.json = kunci — simpan aman, jangan pernah ikut ke AI.")

SUPPORTED = tuple(HANDLERS.keys())
CAT_PREFIX = {cat: pfx for cat, pfx, _, _ in DETECTORS}
CAT_PREFIX["custom"] = "CUSTX"
CATS = list(CAT_PREFIX.keys())


# ---------------------------------------------------------------- helpers

def _save_upload(uploaded):
    suffix = os.path.splitext(uploaded.name)[1].lower()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp.write(uploaded.getvalue())
    tmp.close()
    return tmp.name


def _zip_results(files, extra=None):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data, _ in files:
            if data is not None:
                z.writestr(name, data)
        if extra:
            for name, data in extra.items():
                z.writestr(name, data)
    buf.seek(0)
    return buf.getvalue()


def codebook_from_json_str(s):
    """Codebook from pasted mapping.json. Accepts full file, entries-only, or simple dict."""
    cb = Codebook()
    try:
        data = json.loads(s)
    except json.JSONDecodeError as e:
        return None, f"JSON tidak valid: {e}"
    entries = data.get("entries", data)
    for orig, entry in entries.items():
        if isinstance(entry, dict) and "token" in entry:
            token, cat = entry["token"], entry.get("category", "?")
        elif isinstance(entry, str):
            token, cat = entry, "?"
        else:
            continue
        cb.map[orig] = token
        cb.rev[token] = orig
        cb.categories[orig] = cat
        cb._bump_counter(token)
    if not cb.rev:
        return None, "Tidak ada entri token ditemukan di JSON."
    return cb, None


def mapping_input(label, key):
    """Upload-or-paste widget for mapping.json → Codebook or None."""
    m = st.radio(f"{label} — sumber", ["📁 Upload", "📋 Paste JSON"],
                 horizontal=True, key=f"{key}_src")
    if m == "📁 Upload":
        up = st.file_uploader(f"{label} (mapping.json)", type=["json"], key=f"{key}_up")
        if up:
            p = _save_upload(up)
            cb = Codebook(p)
            os.unlink(p)
            return cb
        return None
    pasted = st.text_area(f"{label} — paste isi mapping.json", height=140, key=f"{key}_txt")
    if pasted.strip():
        cb, err = codebook_from_json_str(pasted)
        if err:
            st.error(err)
            return None
        st.success(f"Codebook loaded: {len(cb.map)} entri")
        return cb
    return None


def decode_table(text, cb):
    rows = []
    for token, orig in cb.rev.items():
        n = text.count(token)
        if n:
            rows.append({"Token": token, "Decoded to": orig,
                         "Category": cb.categories.get(orig, "?"), "Count": n})
    return sorted(rows, key=lambda r: r["Token"])


def codebook_bytes(cb):
    entries = {o: {"token": t, "category": cb.categories.get(o, "?"),
                   "recommendation": MASK_RECOMMENDATION.get(cb.categories.get(o, ""), "")}
               for o, t in cb.map.items()}
    return json.dumps({"meta": cb.meta, "entries": entries}, indent=2).encode()


tab_enc, tab_dec = st.tabs(["🎭 Encode", "🔓 Decode"])

# ================================================================= ENCODE
with tab_enc:
    st.subheader("1 · Input")
    mode = st.radio("Input", ["📁 Upload file", "📋 Paste text"], horizontal=True,
                    key="mask_mode", label_visibility="collapsed")
    with st.expander("(Opsional) mapping.json lama — biar token konsisten antar batch"):
        prev_cb = mapping_input("Mapping lama", "mask_map")

    pasted, ups = None, None
    if mode == "📋 Paste text":
        pasted = st.text_area("Paste teks yang mau di-encode", height=180, key="mask_txt")
    else:
        ups = st.file_uploader("Upload files", type=[e.lstrip(".") for e in SUPPORTED],
                               accept_multiple_files=True, key="mask_up")

    if st.button("🔍 Scan", type="primary", disabled=not (pasted or ups)):
        sources, structured = [], []
        if pasted:
            sources.append(("(pasted text)", pasted))
        if ups:
            for up in ups:
                p = _save_upload(up)
                sources.append((up.name, extract_text(p)))
                if os.path.splitext(up.name)[1].lower() in (".xlsx", ".xlsm", ".csv", ".tsv"):
                    structured.extend((cat, val, col, up.name)
                                      for cat, val, col in find_structured(p))
                os.unlink(p)
        combined = "\n".join(t for _, t in sources)
        st.session_state["mask_scan_text"] = combined
        seen, where = {}, {}
        # column-aware findings first (header context beats regex guess)
        for cat, val, col, fname in structured:
            seen.setdefault(val, cat)
            where.setdefault(val, set()).add(f"{fname} [kolom: {col}]")
        for name, t in sources:
            for cat, pfx, val in find_all(t):
                seen.setdefault(val, cat)
                where.setdefault(val, set()).add(name)
        def fmt_files(names):
            # if "file.xlsx [kolom: X]" exists, drop the bare "file.xlsx" duplicate
            names = set(names)
            bare_dupes = {n for n in names
                          if any(o != n and o.startswith(n + " [") for o in names)}
            return ", ".join(sorted(names - bare_dupes))
        st.session_state["mask_rows"] = [
            {"Encode": True, "Value": v, "Category": c,
             "Matches": combined.count(v), "Files": fmt_files(where.get(v, []))}
            for v, c in sorted(seen.items())
        ]
        st.session_state["mask_ver"] = st.session_state.get("mask_ver", 0) + 1
        st.session_state.pop("enc_results", None)  # new scan invalidates old results

    # ---------------- review & select
    if st.session_state.get("mask_rows"):
        st.subheader("2 · Review & pilih")
        rows = st.session_state["mask_rows"]

        all_cats = sorted({r["Category"] for r in rows})
        cat_counts = {c: sum(1 for r in rows if r["Category"] == c) for c in all_cats}
        c1, c2 = st.columns([2, 1])
        inc_cats = c1.multiselect(
            "Kategori", options=all_cats, default=all_cats,
            format_func=lambda c: f"{c} ({cat_counts[c]})", key="mask_cats")
        q = c2.text_input("🔎 Filter value (substring)", key="mask_filter")

        def visible(r):
            return r["Category"] in inc_cats and (not q or q.lower() in r["Value"].lower())
        vis_idx = [k for k, r in enumerate(rows) if visible(r)]

        b1, b2, b3 = st.columns(3)
        if b1.button(f"✅ Check hasil filter ({len(vis_idx)})"):
            for k in vis_idx:
                rows[k]["Encode"] = True
            st.session_state["mask_ver"] += 1
            st.rerun()
        if b2.button(f"⬜ Uncheck hasil filter ({len(vis_idx)})"):
            for k in vis_idx:
                rows[k]["Encode"] = False
            st.session_state["mask_ver"] += 1
            st.rerun()
        if b3.button(f"🗑 Hapus hasil filter ({len(vis_idx)})"):
            st.session_state["mask_rows"] = [r for k, r in enumerate(rows)
                                             if k not in set(vis_idx)]
            st.session_state["mask_ver"] += 1
            st.rerun()

        vis_rows = [dict(rows[k]) for k in vis_idx]  # copies, editor owns them
        ver = st.session_state.get("mask_ver", 0)
        edited = st.data_editor(
            vis_rows,
            column_config={
                "Encode": st.column_config.CheckboxColumn("Encode", default=True),
                "Value": st.column_config.TextColumn("Value", required=True),
                "Category": st.column_config.SelectboxColumn("Category", options=CATS,
                                                             default="custom"),
                "Matches": st.column_config.NumberColumn("Matches", disabled=True),
                "Files": st.column_config.TextColumn("Ada di file", disabled=True),
            },
            num_rows="dynamic", use_container_width=True,
            key=f"mask_editor_v{ver}_{q}_{','.join(sorted(inc_cats))}",
        )

        # ---- merge back by VALUE (robust against row deletion/reordering)
        by_val = {r["Value"]: r for r in rows}
        edited_vals = set()
        structural_change = False
        for erow in edited:
            v = erow.get("Value")
            if not v:
                continue
            edited_vals.add(v)
            if v in by_val:
                m = by_val[v]
                if (m["Encode"] != bool(erow.get("Encode"))
                        or m["Category"] != (erow.get("Category") or "custom")):
                    m["Encode"] = bool(erow.get("Encode"))
                    m["Category"] = erow.get("Category") or "custom"
            else:  # brand-new / renamed row
                rows.append({"Encode": bool(erow.get("Encode", True)), "Value": v,
                             "Category": erow.get("Category") or "custom",
                             "Matches": st.session_state.get("mask_scan_text", "").count(v),
                             "Files": "(manual)"})
                structural_change = True
        # rows visible before but now missing from editor -> user deleted them
        deleted = [k for k in vis_idx if rows[k]["Value"] not in edited_vals]
        if deleted:
            st.session_state["mask_rows"] = [r for k, r in enumerate(rows)
                                             if k not in set(deleted)]
            structural_change = True
        if structural_change:
            st.session_state["mask_ver"] += 1
            st.rerun()
        rows = st.session_state["mask_rows"]

        selected = {r["Value"]: r.get("Category") or "custom"
                    for r in rows if r.get("Encode") and r.get("Value")
                    and r["Category"] in inc_cats}
        excl = sum(1 for r in rows if r["Category"] not in inc_cats)
        st.write(f"Terpilih untuk encode: **{len(selected)}** dari {len(rows)} value"
                 + (f" · 🚫 {excl} di-exclude via kategori" if excl else ""))

        with st.expander("🔍 Cek konteks — value ini muncul di mana?"):
            ctx_q = st.text_input("Value yang mau dicek:", key="ctx_q")
            scan_text = st.session_state.get("mask_scan_text", "")
            if ctx_q and scan_text:
                hits = [m.start() for m in re.finditer(re.escape(ctx_q), scan_text)]
                if not hits:
                    st.warning("0 match di teks hasil scan.")
                else:
                    st.write(f"{len(hits)} match — menampilkan max 15:")
                    for h in hits[:15]:
                        snip = scan_text[max(0, h - 70):h + len(ctx_q) + 70].replace("\n", " ⏎ ")
                        st.code(f"…{snip}…", language=None)

        # ---------------- encode
        st.subheader("3 · Encode")
        if selected and st.button("🎭 Encode selected", type="primary"):
            cb = prev_cb if prev_cb else Codebook()

            def _alt(v):
                """Escape value + word boundaries so short fragments can't
                corrupt longer values (172.1 must not hit 172.114.11.2)."""
                pre = r"(?<![\w.])" if v[0].isalnum() else ""
                if v[-1].isalnum():
                    suf = r"(?![\w.])" if v[-1].isdigit() else r"(?!\w)"
                elif v[-1] == ".":
                    suf = r"(?!\d)"
                else:
                    suf = ""
                return pre + re.escape(v) + suf

            pattern = re.compile("|".join(
                _alt(v) for v in sorted(selected, key=len, reverse=True)))

            def fn(text):
                cnt = [0]
                def rep(m):
                    cnt[0] += 1
                    orig = m.group(0)
                    cat = selected[orig]
                    return cb.token_for(orig, cat, CAT_PREFIX.get(cat, "CUSTX"))
                return pattern.sub(rep, text), cnt[0]

            files, masked_text, log = [], None, []
            with st.status("🎭 Encoding…", expanded=True) as status:
                st.write(f"▶️ {len(selected)} value terpilih.")
                if pasted:
                    masked_text, n = fn(pasted)
                    log.append(f"pasted text — {n} replacements")
                    st.write(f"✅ {log[-1]}")
                if ups:
                    prog = st.progress(0.0)
                    for i, up in enumerate(ups, 1):
                        ext = os.path.splitext(up.name)[1].lower()
                        handler = HANDLERS.get(ext)
                        if not handler:
                            st.write(f"⏭ {up.name}: tipe {ext} tidak didukung")
                            continue
                        st.write(f"⏳ ({i}/{len(ups)}) {up.name}…")
                        in_path = _save_upload(up)
                        out_path = in_path + ".out" + ext
                        n = handler(in_path, out_path, fn)
                        with open(out_path, "rb") as f:
                            data = f.read()
                        os.unlink(in_path); os.unlink(out_path)
                        stem, e = os.path.splitext(up.name)
                        files.append((f"{stem}_masked{e}", data, n))
                        log.append(f"{up.name} — {n} replacements")
                        st.write(f"✅ ({i}/{len(ups)}) {log[-1]}")
                        prog.progress(i / len(ups))
                status.update(label=f"✅ Encode selesai — {len(cb.map)} unique di codebook",
                              state="complete")

            st.session_state["enc_results"] = {
                "files": files, "masked_text": masked_text, "log": log,
                "map_bytes": codebook_bytes(cb),
                "cb_rows": [{"Original": o, "Token": t,
                             "Category": cb.categories.get(o, "?")}
                            for o, t in cb.map.items()],
            }

    # ---------------- results (persist across reruns / downloads)
    res = st.session_state.get("enc_results")
    if res:
        st.divider()
        st.subheader("📦 Hasil encode")
        for line in res["log"]:
            st.write(f"✅ {line}")
        if res["masked_text"] is not None:
            st.text_area("Hasil masked — copy dari sini, kasih ke AI",
                         res["masked_text"], height=180, key="mask_txt_out")
        st.dataframe(res["cb_rows"], use_container_width=True)
        st.error("⚠️ Simpan mapping.json (download atau copy) — itu kunci decode, "
                 "jangan ikut ke AI.")
        with st.expander("📋 Copy mapping.json sebagai teks"):
            st.code(res["map_bytes"].decode(), language="json")
        d1, d2 = st.columns(2)
        if res["files"]:
            d1.download_button("⬇️ Masked files + mapping.json (zip)",
                               _zip_results(res["files"],
                                            extra={"mapping.json": res["map_bytes"]}),
                               "masked_bundle.zip", "application/zip",
                               type="primary", key="dl_zip")
        d2.download_button("⬇️ mapping.json", res["map_bytes"],
                           "mapping.json", "application/json", key="dl_map")
        for name, data, n in res["files"]:
            st.download_button(f"⬇️ {name}", data, name, key=f"dl_{name}")

# ================================================================= DECODE
with tab_dec:
    st.subheader("1 · Codebook")
    dec_cb = mapping_input("Codebook", "dec_map")

    st.subheader("2 · Input")
    mode = st.radio("Input", ["📁 Upload file", "📋 Paste text"], horizontal=True,
                    key="dec_mode", label_visibility="collapsed")
    d_pasted, d_ups = None, None
    if mode == "📋 Paste text":
        d_pasted = st.text_area("Paste output AI yang berisi token", height=180,
                                key="dec_txt")
    else:
        d_ups = st.file_uploader("Upload files (pptx/docx/xlsx/csv/...)",
                                 type=[e.lstrip(".") for e in SUPPORTED],
                                 accept_multiple_files=True, key="dec_up")

    if st.button("🔓 Decode", type="primary",
                 disabled=not (dec_cb and (d_pasted or d_ups))):
        cb = dec_cb
        files, tables, text_out, total = [], {}, None, 0
        with st.status("🔓 Decoding…", expanded=True) as status:
            if d_pasted:
                tables["(pasted text)"] = decode_table(d_pasted, cb)
                text_out, n = unmask_text(d_pasted, cb)
                total += n
                st.write(f"✅ pasted text — {n} token direstorasi")
            if d_ups:
                prog = st.progress(0.0)
                for i, up in enumerate(d_ups, 1):
                    ext = os.path.splitext(up.name)[1].lower()
                    handler = HANDLERS.get(ext)
                    if not handler:
                        st.write(f"⏭ {up.name}: tipe {ext} tidak didukung")
                        continue
                    st.write(f"⏳ ({i}/{len(d_ups)}) {up.name}…")
                    in_path = _save_upload(up)
                    tables[up.name] = decode_table(extract_text(in_path), cb)
                    out_path = in_path + ".out" + ext
                    n = handler(in_path, out_path, lambda t: unmask_text(t, cb))
                    with open(out_path, "rb") as f:
                        data = f.read()
                    os.unlink(in_path); os.unlink(out_path)
                    stem, e = os.path.splitext(up.name)
                    files.append((f"{stem}_restored{e}", data, n))
                    total += n
                    st.write(f"✅ ({i}/{len(d_ups)}) {up.name} — {n} token direstorasi")
                    prog.progress(i / len(d_ups))
            status.update(label=f"✅ Decode selesai — {total} token direstorasi",
                          state="complete")
        st.session_state["dec_results"] = {"files": files, "tables": tables,
                                           "text": text_out, "total": total}

    res = st.session_state.get("dec_results")
    if res:
        st.divider()
        st.subheader("📦 Hasil decode")
        if res["total"] == 0:
            st.info("Tidak ada token ditemukan — cek apakah mapping.json-nya benar "
                    "untuk batch ini, atau AI mengubah format token.")
        for src_name, table in res["tables"].items():
            if table:
                st.markdown(f"**Decode summary — {src_name}**")
                st.dataframe(table, use_container_width=True)
        if res["text"] is not None:
            st.text_area("Hasil decoded", res["text"], height=180, key="dec_txt_out")
        for name, data, n in res["files"]:
            st.download_button(f"⬇️ {name} ({n} restorasi)", data, name,
                               key=f"dl_dec_{name}")
        if len(res["files"]) > 1:
            st.download_button("⬇️ Download semua (zip)", _zip_results(res["files"]),
                               "restored_bundle.zip", "application/zip",
                               type="primary", key="dl_dec_zip")
