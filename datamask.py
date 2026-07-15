#!/usr/bin/env python3
"""
datamask.py — Reversible sensitive-data masking for SOC/cybersecurity reports.

Workflow:
  1) scan    : detect sensitive data (IP, hostname, username, email, MAC, domain)
               across .docx / .xlsx / .csv / .txt / .pptx files
  2) mask    : replace findings with deterministic tokens (IPX001, HOSTX001, ...)
               and write a codebook (mapping.json) for later restoration
  3) unmask  : restore original values in ANY supported file (e.g. an AI-generated
               .pptx built from masked data) using the codebook

Token format is a single alphanumeric word (e.g. HOSTX003) so LLMs and Office
tools carry it through untouched and it survives copy/paste, tables, charts, etc.

Usage:
  python3 datamask.py scan   report.docx data.xlsx events.csv
  python3 datamask.py mask   report.docx data.xlsx events.csv -o masked/ -m mapping.json
  python3 datamask.py unmask ai_generated.pptx -m mapping.json -o restored/
  python3 datamask.py report -m mapping.json          # masking recommendation list

The codebook (mapping.json) is the secret — protect it like a key.
"""

import argparse
import csv as csv_mod
import io
import json
import os
import re
import sys
from collections import OrderedDict
from datetime import datetime, timezone

# ---------------------------------------------------------------- detectors

PRIVATE_NETS = re.compile(
    r"^(10\.|192\.168\.|172\.(1[6-9]|2[0-9]|3[01])\.|127\.|169\.254\.)"
)

def _valid_ipv4(s: str) -> bool:
    parts = s.split(".")
    return len(parts) == 4 and all(p.isdigit() and 0 <= int(p) <= 255 for p in parts)

# Order matters: more specific first (email before domain, FQDN before hostname).
DETECTORS = [
    # (category, token prefix, regex, post-filter)
    ("email",    "MAILX", re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"), None),
    ("ipv4",     "IPX",   re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b"), _valid_ipv4),
    ("ipv6",     "IP6X",  re.compile(r"\b(?:[0-9A-Fa-f]{1,4}:){2,7}[0-9A-Fa-f]{1,4}\b"),
        lambda s: not re.fullmatch(r"(?:[0-9A-Fa-f]{2}:){5}[0-9A-Fa-f]{2}", s)),
    ("mac",      "MACX",  re.compile(r"\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b"), None),
    # Windows-style hostnames: DESKTOP-XXXX, SRV-*, WIN-*, DC01, LAPTOP-...
    ("hostname", "HOSTX", re.compile(
        r"\b(?:DESKTOP|LAPTOP|WIN|SRV|SVR|DC|PC|WS|HOST|VM|APP|DB|WEB|MAIL|FW|SW|RTR|PRD|DEV|UAT)"
        r"[-_][A-Za-z0-9][A-Za-z0-9-_]{1,30}\b", re.IGNORECASE), None),
    # FQDN (internal domains etc.) — post-filter drops common file extensions
    ("fqdn",     "FQDNX", re.compile(
        r"\b(?:[A-Za-z0-9](?:[A-Za-z0-9-]{0,61}[A-Za-z0-9])?\.){2,}[A-Za-z]{2,}\b"),
        lambda s: not re.search(r"\.(?:exe|dll|csv|xlsx|docx|pptx|txt|log|json|xml|py|sh|js|zip|png|jpg)$", s, re.I)),
    # DOMAIN\username
    ("username", "USERX", re.compile(r"\b[A-Za-z][A-Za-z0-9_-]{1,20}\\[A-Za-z][A-Za-z0-9._-]{2,30}\b"), None),
    # key=value contextual usernames: user=jdoe, username: budi.s, account="ops_admin"
    ("username", "USERX", re.compile(
        r"(?i)\b(?:user(?:name)?|account|login|logon|acct)\s*[:=]\s*\"?([A-Za-z][A-Za-z0-9._-]{2,30})\"?"), None),
]

MASK_RECOMMENDATION = {
    "ipv4":     "Tokenize (reversible). If sharing externally without codebook: mask last octet (10.1.2.x).",
    "ipv6":     "Tokenize (reversible). External: truncate to /48 prefix.",
    "hostname": "Tokenize (reversible). External: role-based alias (WEB-SRV-A).",
    "fqdn":     "Tokenize (reversible). External: replace internal domain with example.internal.",
    "username": "Tokenize (reversible). Never share real usernames outside the org — high phishing value.",
    "email":    "Tokenize (reversible). External: hash local-part, keep domain if public.",
    "mac":      "Tokenize (reversible). External: keep OUI (vendor) bytes, mask last 3 octets.",
}

# ---------------------------------------------------------------- codebook

class Codebook:
    def __init__(self, path=None):
        self.map = OrderedDict()      # original -> token
        self.rev = OrderedDict()      # token -> original
        self.meta = {"created": datetime.now(timezone.utc).isoformat(), "tool": "datamask v1.0"}
        self.counters = {}
        if path and os.path.exists(path):
            with open(path) as f:
                data = json.load(f)
            self.meta = data.get("meta", self.meta)
            for orig, entry in data.get("entries", {}).items():
                self.map[orig] = entry["token"]
                self.rev[entry["token"]] = orig
                self._bump_counter(entry["token"])
            self.categories = {o: e["category"] for o, e in data.get("entries", {}).items()}
        else:
            self.categories = {}

    def _bump_counter(self, token):
        m = re.match(r"([A-Z0-9]+?X)(\d+)$", token)
        if m:
            pfx, n = m.group(1), int(m.group(2))
            self.counters[pfx] = max(self.counters.get(pfx, 0), n)

    def token_for(self, original, category, prefix):
        if original in self.map:
            return self.map[original]
        self.counters[prefix] = self.counters.get(prefix, 0) + 1
        token = f"{prefix}{self.counters[prefix]:03d}"
        self.map[original] = token
        self.rev[token] = original
        self.categories[original] = category
        return token

    def save(self, path):
        entries = {o: {"token": t, "category": self.categories.get(o, "?"),
                       "recommendation": MASK_RECOMMENDATION.get(self.categories.get(o, ""), "Tokenize (reversible).")}
                   for o, t in self.map.items()}
        with open(path, "w") as f:
            json.dump({"meta": self.meta, "entries": entries}, f, indent=2)
        try:
            os.chmod(path, 0o600)
        except OSError:
            pass

# ---------------------------------------------------------------- core replace

def find_all(text):
    """Yield (category, prefix, match_text) for every detection in text."""
    out = []
    for cat, pfx, rx, flt in DETECTORS:
        for m in rx.finditer(text):
            val = m.group(1) if m.groups() else m.group(0)
            if flt and not flt(val):
                continue
            out.append((cat, pfx, val))
    return out

def mask_text(text, cb: Codebook):
    """Replace all detections in text with tokens. Longest-first to avoid partial overlap."""
    findings = find_all(text)
    if not findings:
        return text, 0
    # Assign tokens
    for cat, pfx, val in findings:
        cb.token_for(val, cat, pfx)
    # Replace longest originals first
    n = 0
    for orig in sorted({v for _, _, v in findings}, key=len, reverse=True):
        token = cb.map[orig]
        new = text.replace(orig, token)
        if new != text:
            n += text.count(orig)
            text = new
    return text, n

def unmask_text(text, cb: Codebook):
    n = 0
    for token in sorted(cb.rev, key=len, reverse=True):
        if token in text:
            n += text.count(token)
            text = text.replace(token, cb.rev[token])
    return text, n

# ---------------------------------------------------------------- file handlers

def _iter_docx_paragraphs(doc):
    """Yield every paragraph in body, tables (nested), headers, footers."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    def walk(parent):
        for p in parent.paragraphs:
            yield p
        for t in parent.tables:
            for row in t.rows:
                for cell in row.cells:
                    yield from walk(cell)
    yield from walk(doc)
    for sec in doc.sections:
        for part in (sec.header, sec.footer):
            yield from walk(part)

def _replace_in_paragraph(par, fn):
    """Apply fn(text)->(text,count) to full paragraph text, preserving run 0 formatting
    when a match spans runs. Only rewrites the paragraph if something changed."""
    full = par.text
    new, n = fn(full)
    if n == 0:
        return 0
    # Try run-by-run first (keeps formatting perfectly when matches are within a run)
    total_in_runs = 0
    for run in par.runs:
        rnew, rn = fn(run.text)
        if rn:
            run.text = rnew
            total_in_runs += rn
    if par.text == new:
        return n
    # Cross-run match remains: collapse into first run
    if par.runs:
        par.runs[0].text = new
        for run in par.runs[1:]:
            run.text = ""
    return n

def process_docx(path, out_path, fn):
    import docx
    doc = docx.Document(path)
    count = 0
    for par in _iter_docx_paragraphs(doc):
        count += _replace_in_paragraph(par, fn)
    doc.save(out_path)
    return count

def process_pptx(path, out_path, fn):
    from pptx import Presentation
    prs = Presentation(path)
    count = 0
    def handle_tf(tf):
        nonlocal count
        for par in tf.paragraphs:
            full = "".join(r.text for r in par.runs)
            new, n = fn(full)
            if n == 0:
                continue
            done_in_runs = 0
            for r in par.runs:
                rnew, rn = fn(r.text)
                if rn:
                    r.text = rnew
                    done_in_runs += rn
            cur = "".join(r.text for r in par.runs)
            if cur != new and par.runs:
                par.runs[0].text = new
                for r in par.runs[1:]:
                    r.text = ""
            count += n
    def handle_shape(shape):
        if shape.has_text_frame:
            handle_tf(shape.text_frame)
        if shape.shape_type == 6:  # group
            for s in shape.shapes:
                handle_shape(s)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    handle_tf(cell.text_frame)
    for slide in prs.slides:
        for shape in slide.shapes:
            handle_shape(shape)
        if slide.has_notes_slide:
            handle_tf(slide.notes_slide.notes_text_frame)
    prs.save(out_path)
    return count

def process_xlsx(path, out_path, fn):
    import openpyxl
    wb = openpyxl.load_workbook(path)  # formulas preserved; we only touch strings
    count = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and not cell.value.startswith("="):
                    new, n = fn(cell.value)
                    if n:
                        cell.value = new
                        count += n
        # sheet name itself can contain hostnames
        new_title, n = fn(ws.title)
        if n:
            ws.title = new_title[:31]
            count += n
    wb.save(out_path)
    return count

def process_textlike(path, out_path, fn):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    new, n = fn(text)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(new)
    return n

HANDLERS = {
    ".docx": process_docx,
    ".pptx": process_pptx,
    ".xlsx": process_xlsx,
    ".xlsm": process_xlsx,
    ".csv":  process_textlike,
    ".tsv":  process_textlike,
    ".txt":  process_textlike,
    ".log":  process_textlike,
    ".json": process_textlike,
    ".xml":  process_textlike,
    ".md":   process_textlike,
}

def extract_text(path):
    """Read all visible text from a file for scan-only mode."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        import docx
        doc = docx.Document(path)
        return "\n".join(p.text for p in _iter_docx_paragraphs(doc))
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            chunks.append(cell.text)
        return "\n".join(chunks)
    if ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True)
        chunks = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                chunks.extend(str(v) for v in row if isinstance(v, str))
        return "\n".join(chunks)
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

# ---------------------------------------------------------------- commands

def cmd_scan(args):
    cb = Codebook()
    per_file = {}
    for path in args.files:
        text = extract_text(path)
        findings = find_all(text)
        agg = {}
        for cat, pfx, val in findings:
            agg.setdefault(cat, set()).add(val)
        per_file[path] = agg
    print(f"\n{'='*70}\n DATAMASK SCAN REPORT — {datetime.now():%Y-%m-%d %H:%M}\n{'='*70}")
    grand = {}
    for path, agg in per_file.items():
        print(f"\n[FILE] {path}")
        if not agg:
            print("  (no sensitive data detected)")
        for cat in sorted(agg):
            vals = sorted(agg[cat])
            grand.setdefault(cat, set()).update(vals)
            print(f"  {cat.upper():9s} ({len(vals):3d} unique): {', '.join(vals[:6])}"
                  + (" ..." if len(vals) > 6 else ""))
    print(f"\n{'-'*70}\n RECOMMENDATIONS\n{'-'*70}")
    for cat in sorted(grand):
        print(f"  {cat.upper():9s} [{len(grand[cat])} unique] -> {MASK_RECOMMENDATION.get(cat)}")
    print(f"\nNext step: python3 datamask.py mask <files> -o masked/ -m mapping.json\n")

def _outpath(path, outdir, suffix):
    base = os.path.basename(path)
    stem, ext = os.path.splitext(base)
    return os.path.join(outdir, f"{stem}{suffix}{ext}")

def cmd_mask(args):
    cb = Codebook(args.mapping if args.append else None)
    os.makedirs(args.outdir, exist_ok=True)
    fn = lambda t: mask_text(t, cb)
    for path in args.files:
        ext = os.path.splitext(path)[1].lower()
        handler = HANDLERS.get(ext)
        if not handler:
            print(f"[skip] {path}: unsupported type {ext}")
            continue
        out = _outpath(path, args.outdir, "_masked")
        n = handler(path, out, fn)
        print(f"[mask] {path} -> {out}  ({n} replacements)")
    cb.save(args.mapping)
    print(f"[codebook] {args.mapping}  ({len(cb.map)} unique values)  — KEEP THIS SECRET")

def cmd_unmask(args):
    cb = Codebook(args.mapping)
    if not cb.rev:
        sys.exit(f"error: codebook {args.mapping} empty or missing")
    os.makedirs(args.outdir, exist_ok=True)
    fn = lambda t: unmask_text(t, cb)
    for path in args.files:
        ext = os.path.splitext(path)[1].lower()
        handler = HANDLERS.get(ext)
        if not handler:
            print(f"[skip] {path}: unsupported type {ext}")
            continue
        out = _outpath(path, args.outdir, "_restored")
        n = handler(path, out, fn)
        print(f"[unmask] {path} -> {out}  ({n} restorations)")

def cmd_report(args):
    cb = Codebook(args.mapping)
    print(f"\n{'ORIGINAL':40s} {'TOKEN':12s} {'CATEGORY':10s} RECOMMENDATION")
    print("-" * 110)
    for orig, token in cb.map.items():
        cat = cb.categories.get(orig, "?")
        print(f"{orig:40s} {token:12s} {cat:10s} {MASK_RECOMMENDATION.get(cat, '')[:50]}")
    print(f"\nTotal: {len(cb.map)} unique sensitive values\n")

def main():
    ap = argparse.ArgumentParser(description="Reversible sensitive-data masking for SOC reports")
    sub = ap.add_subparsers(dest="cmd", required=True)

    p = sub.add_parser("scan", help="detect sensitive data, print report")
    p.add_argument("files", nargs="+")
    p.set_defaults(func=cmd_scan)

    p = sub.add_parser("mask", help="tokenize sensitive data, write codebook")
    p.add_argument("files", nargs="+")
    p.add_argument("-o", "--outdir", default="masked")
    p.add_argument("-m", "--mapping", default="mapping.json")
    p.add_argument("--append", action="store_true", help="extend an existing codebook (consistent tokens across batches)")
    p.set_defaults(func=cmd_mask)

    p = sub.add_parser("unmask", help="restore originals in any file using codebook")
    p.add_argument("files", nargs="+")
    p.add_argument("-o", "--outdir", default="restored")
    p.add_argument("-m", "--mapping", default="mapping.json")
    p.set_defaults(func=cmd_unmask)

    p = sub.add_parser("report", help="print codebook as recommendation list")
    p.add_argument("-m", "--mapping", default="mapping.json")
    p.set_defaults(func=cmd_report)

    args = ap.parse_args()
    args.func(args)

if __name__ == "__main__":
    main()
